import os
import json
import re
import requests
import time
from typing import Dict, List, Any, Optional, Union
from pathlib import Path
import docx
import pdfplumber
from pptx import Presentation
from langchain_core.tools import tool
from .decorators import log_io
from typing import Annotated
from pptx.util import Pt
from dotenv import load_dotenv
from langchain_core.messages import HumanMessage
from src.utils.db_manager import log_handler, ppt_state
import logging
from src.utils.msg_process import resolve_file_path

logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)
logger.addHandler(log_handler)

# sys.path.insert(0, str(Path(__file__).parent))
from src.tools.web_search import web_search_ppt
load_dotenv()
current_index = os.getcwd()
# 假设这些是全局变量
DEEPSEEK_API_KEY = os.getenv('DS_API_KEY')
DEEPSEEK_BASE_URL = os.getenv("DS_API_BASE")
TEMPLATE_LIBRARY_PATH = f'{current_index}/src/tools/templates'
DEFAULT_TEMPLATE_PATH = f'{current_index}/src/tools/templates/deep_green_template.pptx'



    
def call_deepseek_api(prompt: str, max_tokens: int = 2000) -> str:
    """调用DeepSeek API"""
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}"
    }

    data = {
        "model": "deepseek-chat",
        "messages": [
            {
                "role": "user",
                "content": prompt
            }
        ],
        "max_tokens": max_tokens,
        "temperature": 0.7
    }

    try:
        response = requests.post(
            f"{DEEPSEEK_BASE_URL}/chat/completions",
            headers=headers,
            json=data,
            timeout=60
        )

        if response.status_code == 200:
            result = response.json()
            return result["choices"][0]["message"]["content"]
        else:
            print(f"API调用失败: {response.status_code} - {response.text}")
            return None

    except Exception as e:
        print(f"调用DeepSeek API时出错: {e}")
        return None




def load_template_descriptions() -> List[Dict[str, str]]:
    """
    加载模板库描述

    Returns:
        模板描述列表
    """
    description_file = f'{current_index}/src/tools/templates/description.txt'
    templates = []

    try:
        with open(description_file, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        for line in lines:
            line = line.strip()
            if not line or ':' not in line:
                continue

            parts = line.split(':', 2)
            if len(parts) >= 3:
                template_name = parts[0].strip()
                template_path = parts[1].strip()
                template_description = parts[2].strip()

                # 确保路径是绝对路径
                if not os.path.isabs(template_path):
                    template_path = os.path.join(TEMPLATE_LIBRARY_PATH, template_path)

                templates.append({
                    "name": template_name,
                    "path": template_path,
                    "description": template_description
                })
    except Exception as e:
        print(f"加载模板描述失败: {e}")

    return templates


def select_template_by_topic(topic: str) -> str:
    """
    根据主题自动选择PPT模板

    Args:
        topic: PPT主题

    Returns:
        选中的模板路径
    """
    templates = load_template_descriptions()

    if not templates:
        logger.info(f"{ppt_state}-=-PPT专员===未找到可用模板，使用默认模板")
        return DEFAULT_TEMPLATE_PATH

    # 构建模板选择提示
    template_options = "\n".join([
        f"- {t['name']}: {t['description']} (路径: {t['path']})"
        for t in templates
    ])

    prompt = f"""
    请根据以下PPT主题，从可用的PPT模板中选择最合适的一个：

    PPT主题: {topic}

    可用模板:
    {template_options}

    请根据主题内容、风格和适用场景选择最匹配的模板。
    只返回你选择的模板名称，不要包含其他任何内容。
    """

    selected_template_name = call_deepseek_api(prompt, max_tokens=100)

    if not selected_template_name:
        print("模板选择失败，使用默认模板")
        return DEFAULT_TEMPLATE_PATH

    # 清理返回的模板名称
    selected_template_name = selected_template_name.strip().replace('"', '').replace("'", "")

    # 查找匹配的模板
    for template in templates:
        if template["name"] == selected_template_name:
            logger.info(f"{ppt_state}-=-PPT专员===已选择模板: {template['name']} - {template['path']}")
            return template["path"]

    # 如果没有找到匹配的模板，使用默认模板
    logger.info(f"{ppt_state}-=-PPT专员===未找到模板 '{selected_template_name}'，使用默认模板")
    return DEFAULT_TEMPLATE_PATH


def analyze_input_type(user_input: str) -> Dict[str, Any]:
    """
    分析用户输入类型

    Args:
        user_input: 用户输入

    Returns:
        分析结果
    """
    prompt = f"""
    请分析以下用户输入，判断其类型和内容完整性：

    用户输入: {user_input}

    请判断：
    1. 这是一个PPT主题、PPT制作需求、文档内容，还是PPT大纲？
    2. 内容是否完整，是否需要通过网页搜索补充信息？

    输出格式为JSON：
    {{
        "input_type": "theme|requirement|document|outline",
        "is_complete": true|false,
        "extracted_topic": "提取的主题（如果有）",
        "needs_web_search": true|false,
        "search_query": "需要搜索的查询词"
    }}
    """

    response = call_deepseek_api(prompt)
    if not response:
        return {
            "input_type": "unknown",
            "is_complete": False,
            "extracted_topic": "",
            "needs_web_search": True,
            "search_query": user_input
        }

    try:
        json_match = re.search(r'\{.*\}', response, re.DOTALL)
        if json_match:
            return json.loads(json_match.group())
        else:
            return {
                "input_type": "unknown",
                "is_complete": False,
                "extracted_topic": "",
                "needs_web_search": True,
                "search_query": user_input
            }
    except json.JSONDecodeError:
        return {
            "input_type": "unknown",
            "is_complete": False,
            "extracted_topic": "",
            "needs_web_search": True,
            "search_query": user_input
        }


def extract_topic_from_content(content: str) -> str:
    """
    从内容中提取主题

    Args:
        content: 内容文本

    Returns:
        提取的主题
    """
    prompt = f"""
    请从以下内容中提取一个合适的PPT主题：

    {content[:1000]}

    请输出一个简洁明确的主题，不要包含任何其他内容。
    """

    response = call_deepseek_api(prompt, max_tokens=100)
    return response.strip() if response else "未知主题"


def load_document_content(document_path: str) -> str:
    """
    从文档中加载内容

    Args:
        document_path: 文档路径

    Returns:
        文档内容
    """
    file_extension = Path(document_path).suffix.lower()

    try:
        if file_extension == '.docx':
            return load_word_document(document_path)
        elif file_extension == '.pdf':
            return load_pdf_document(document_path)
        elif file_extension == '.md':
            return load_markdown_document(document_path)
        else:
            return f"不支持的文件格式: {file_extension}"
    except Exception as e:
        return f"加载文档失败: {e}"


def load_word_document(docx_path: str) -> str:
    """加载Word文档"""
    doc = docx.Document(docx_path)
    full_text = []

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            full_text.append(paragraph.text)

    return '\n'.join(full_text)


def load_pdf_document(pdf_path: str) -> str:
    """加载PDF文档"""
    full_text = []

    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                full_text.append(text)

    return '\n'.join(full_text)


def load_markdown_document(md_path: str) -> str:
    """加载Markdown文档"""
    try:
        with open(md_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # 预处理Markdown内容
        return preprocess_markdown(content)
    except Exception as e:
        return f"读取Markdown文件失败: {e}"


def preprocess_markdown(content: str) -> str:
    """
    预处理Markdown内容

    Args:
        content: 原始Markdown内容

    Returns:
        处理后的文本内容
    """
    # 移除代码块
    content = re.sub(r'```.*?```', '', content, flags=re.DOTALL)

    # 移除行内代码标记
    content = re.sub(r'`([^`]+)`', r'\1', content)

    # 保留标题标记但简化格式
    content = re.sub(r'#+\s*(.+)', r'标题: \1', content)

    # 处理列表项
    content = re.sub(r'^\s*[-*+]\s+(.+)', r'• \1', content, flags=re.MULTILINE)
    content = re.sub(r'^\s*\d+\.\s+(.+)', r'\1', content, flags=re.MULTILINE)

    # 移除链接和图片标记
    content = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', content)
    content = re.sub(r'!\[([^\]]*)\]\([^)]+\)', r'', content)

    # 移除粗体和斜体标记
    content = re.sub(r'\*\*(.+?)\*\*', r'\1', content)
    content = re.sub(r'\*(.+?)\*', r'\1', content)
    content = re.sub(r'_(.+?)_', r'\1', content)

    # 移除多余的空行
    content = re.sub(r'\n\s*\n', '\n\n', content)

    return content.strip()


def enhance_content_with_web_search(topic: str, original_content: str = "") -> str:
    """
    使用网页搜索增强内容

    Args:
        topic: 主题
        original_content: 原始内容

    Returns:
        增强后的内容
    """
    # 进行网页搜索
    search_results = web_search_ppt(topic)
    # search_results = search['content']
    # if not search_results:
    #     return original_content

    # # 使用大模型整合原始内容和搜索内容
    # prompt = f"""
    # 请基于以下原始内容和网页搜索结果，为PPT生成更全面、准确的内容：

    # 主题: {topic}

    # 原始内容:
    # {original_content}

    # 网页搜索结果:
    # {search_results}

    # 请整合这些信息，生成适合PPT使用的专业内容，突出重点和关键信息。
    # 输出整合后的内容。
    # """

    # enhanced_content = call_deepseek_api(prompt, max_tokens=2000)
    return search_results


def generate_ppt_structure(topic: str, content: str, template_slide_count: int) -> Dict[str, Any]:
    """
    生成PPT结构

    Args:
        topic: PPT主题
        content: 内容
        template_slide_count: 模板幻灯片数量

    Returns:
        PPT结构
    """
    prompt = f"""
    请为关于"{topic}"的PPT规划整体结构。

    可用内容:
    {content[:3000]}

    模板幻灯片数量: {template_slide_count}

    请设计一个逻辑清晰的PPT内容结构，包含封面、目录、内容页和总结。

    输出格式要求为JSON:
    {{
        "title": "PPT主标题",
        "subtitle": "PPT副标题",
        "slides": [
            {{
                "slide_index": 0,
                "purpose": "封面",
                "title": "封面标题",
                "content": ["封面内容要点1", "封面内容要点2"]
            }},
            {{
                "slide_index": 1,
                "purpose": "目录",
                "title": "目录标题",
                "content": ["目录项1", "目录项2", "目录项3"]
            }},
            // 更多幻灯片...
        ]
    }}
    """

    response = call_deepseek_api(prompt)
    if not response:
        return create_default_structure(topic, template_slide_count)

    try:
        json_match = re.search(r'\{.*\}', response, re.DOTALL)
        if json_match:
            return json.loads(json_match.group())
        else:
            return create_default_structure(topic, template_slide_count)
    except json.JSONDecodeError:
        return create_default_structure(topic, template_slide_count)


def create_default_structure(topic: str, slide_count: int) -> Dict[str, Any]:
    """创建默认PPT结构"""
    structure = {
        "title": topic,
        "subtitle": f"关于{topic}的专题报告",
        "slides": []
    }

    # 添加封面
    structure['slides'].append({
        "slide_index": 0,
        "purpose": "封面",
        "title": topic,
        "content": [f"关于{topic}的专题报告", "演讲人: 待填写", "日期: 待填写"]
    })

    # 添加目录（如果有足够页面）
    if slide_count >= 2:
        structure['slides'].append({
            "slide_index": 1,
            "purpose": "目录",
            "title": "目录",
            "content": [f"{topic}概述", f"{topic}的核心要素", f"{topic}的应用场景", f"{topic}的未来发展", "总结"]
        })

    # 添加内容页
    content_slide_count = max(1, slide_count - 3)  # 减去封面、目录和总结

    for i in range(content_slide_count):
        structure['slides'].append({
            "slide_index": i + 2,
            "purpose": "内容页",
            "title": f"{topic}的内容部分{i + 1}",
            "content": [f"关于{topic}的重要内容点1", f"关于{topic}的重要内容点2", f"关于{topic}的重要内容点3"]
        })

    # 添加总结（如果有足够页面）
    if slide_count >= 3:
        structure['slides'].append({
            "slide_index": slide_count - 1,
            "purpose": "总结",
            "title": "总结",
            "content": [f"{topic}的重要性", "关键要点回顾", "未来展望"]
        })

    return structure


def generate_slide_content(slide_info: Dict[str, Any], topic: str, content: str) -> Dict[str, Any]:
    """
    生成单页幻灯片内容

    Args:
        slide_info: 幻灯片信息
        topic: PPT主题
        content: 可用内容

    Returns:
        幻灯片内容
    """
    slide_purpose = slide_info.get('purpose', '内容页')

    if slide_purpose == '封面':
        prompt = f"""
        请为关于"{topic}"的PPT封面页生成内容。

        可用内容:
        {content[:1000]}

        要求:
        - 主标题: {slide_info.get('title', topic)}
        - 生成2-3个适合封面的内容要点

        输出格式:
        {{
            "title": "封面标题",
            "content": ["要点1", "要点2", "要点3"]
        }}
        """

    elif slide_purpose == '目录':
        prompt = f"""
        请为关于"{topic}"的PPT目录页生成内容。

        可用内容:
        {content[:1000]}

        要求:
        - 标题: {slide_info.get('title', '目录')}
        - 生成5-7个逻辑连贯的目录项

        输出格式:
        {{
            "title": "目录",
            "content": ["目录项1", "目录项2", "目录项3", "..."]
        }}
        """

    elif slide_purpose == '内容页':
        prompt = f"""
        请为关于"{topic}"的PPT内容页生成详细内容。

        可用内容:
        {content[:1500]}

        要求:
        - 标题: {slide_info.get('title', f'{topic}相关内容')}
        - 生成3-5个详细的内容要点，每个要点可以有一些解释
        - 内容要专业、有条理

        输出格式:
        {{
            "title": "页面标题",
            "content": ["要点1及简要解释", "要点2及简要解释", "要点3及简要解释"]
        }}
        """

    elif slide_purpose == '总结':
        prompt = f"""
        请为关于"{topic}"的PPT总结页生成内容。

        可用内容:
        {content[:1000]}

        要求:
        - 标题: {slide_info.get('title', '总结')}
        - 生成3-4个总结要点，回顾核心内容
        - 可以包含未来展望或建议

        输出格式:
        {{
            "title": "总结",
            "content": ["总结要点1", "总结要点2", "总结要点3", "未来展望"]
        }}
        """

    else:
        prompt = f"""
        请为关于"{topic}"的PPT{slide_purpose}页生成内容。

        可用内容:
        {content[:1000]}

        要求:
        - 标题: {slide_info.get('title', slide_purpose)}
        - 生成2-3个适合{slide_purpose}页的内容

        输出格式:
        {{
            "title": "页面标题",
            "content": ["内容1", "内容2"]
        }}
        """

    response = call_deepseek_api(prompt, max_tokens=1000)
    if not response:
        return slide_info

    try:
        json_match = re.search(r'\{.*\}', response, re.DOTALL)
        if json_match:
            slide_content = json.loads(json_match.group())
            slide_content['purpose'] = slide_purpose
            slide_content['slide_index'] = slide_info.get('slide_index', 0)
            return slide_content
        else:
            return slide_info
    except json.JSONDecodeError:
        return slide_info


def replace_ppt_content(template_path: str, ppt_structure: Dict[str, Any], slide_contents: Dict[int, Any],
                        output_path: str):
    """
    替换PPT内容

    Args:
        template_path: 模板路径
        ppt_structure: PPT结构
        slide_contents: 幻灯片内容
        output_path: 输出路径
    """
    try:
        prs = Presentation(template_path)

        # 确保目录项与内容页标题一致
        directory_slide_index = None
        directory_content = None

        for slide_index, content in slide_contents.items():
            if content.get('purpose') == '目录':
                directory_slide_index = slide_index
                directory_content = content.get('content', [])
                break

        if directory_content:
            content_titles = []
            for slide_index, content in slide_contents.items():
                if content.get('purpose') == '内容页':
                    content_titles.append(content.get('title', ''))

            if directory_content != content_titles:
                updated_directory = content_titles[:len(directory_content)] if len(content_titles) < len(
                    directory_content) else content_titles
                while len(updated_directory) < len(directory_content):
                    updated_directory.append(f"内容部分 {len(updated_directory) + 1}")
                slide_contents[directory_slide_index]['content'] = updated_directory

        # 替换每页内容
        for slide_index, content in slide_contents.items():
            if slide_index < len(prs.slides):
                slide = prs.slides[slide_index]

                # 简单替换策略：第一个形状作为标题，其他作为内容
                text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text') and shape.text.strip()]

                if text_shapes and content.get('title'):
                    # 替换标题
                    title_shape = text_shapes[0]
                    if hasattr(title_shape, 'text_frame'):
                        title_shape.text_frame.clear()
                        p = title_shape.text_frame.paragraphs[
                            0] if title_shape.text_frame.paragraphs else title_shape.text_frame.add_paragraph()
                        p.text = content['title']

                if len(text_shapes) > 1 and content.get('content'):
                    # 替换内容
                    content_shape = text_shapes[1]
                    if hasattr(content_shape, 'text_frame'):
                        content_shape.text_frame.clear()
                        p = content_shape.text_frame.paragraphs[
                            0] if content_shape.text_frame.paragraphs else content_shape.text_frame.add_paragraph()
                        content_text = '\n'.join([f"• {item}" for item in content['content']])
                        p.text = content_text

        prs.save(output_path)
        return True
    except Exception as e:
        print(f"替换PPT内容失败: {e}")
        return False


# 主工具函数
def ppt_generation_tool(
        user_input: str,
        document_path: str = None
) -> str:
    """
    PPT生成工具

    Args:
        user_input: 用户输入（主题、需求、大纲等）
        document_path: 文档路径（可选）

    Returns:
        生成结果信息
    """
    # 分析输入类型
    input_analysis = analyze_input_type(user_input)
    input_type = input_analysis.get('input_type', 'unknown')
    needs_web_search = input_analysis.get('needs_web_search', True)
    search_query = input_analysis.get('search_query', user_input)

    # 提取或确定主题
    if input_type in ['document', 'outline']:
        topic = input_analysis.get('extracted_topic', '')
        if not topic:
            # 从文档中提取主题
            # if document_path:
            #     document_content = load_document_content(document_path)
            #     topic = extract_topic_from_content(document_content)
            # else:
                topic = extract_topic_from_content(user_input)
    else:
        topic = input_analysis.get('extracted_topic', user_input)

    if not topic:
        topic = "未知主题"

    # 根据主题自动选择模板
    selected_template_path = select_template_by_topic(topic)

    # 收集内容
    content_parts = []

    # 添加用户输入
    if input_type in ['requirement', 'outline']:
        content_parts.append(f"用户需求/大纲:\n{user_input}")

    # 添加文档内容
    # if document_path:
    #     document_content = load_document_content(document_path)
    #     if not document_content.startswith("不支持") and not document_content.startswith("加载文档失败"):
    #         content_parts.append(f"文档内容:\n{document_content}")

    # 网页搜索补充
    if needs_web_search:
        web_content = enhance_content_with_web_search(search_query, "\n".join(content_parts))
        content_parts.append(f"网页搜索补充:\n{web_content}")

    # 合并所有内容
    full_content = "\n\n".join(content_parts)

    # 加载模板并分析结构
    try:
        prs = Presentation(selected_template_path)
        template_slide_count = len(prs.slides)
    except Exception as e:
        return f"加载PPT模板失败: {e}"

    # 生成PPT结构
    ppt_structure = generate_ppt_structure(topic, full_content, template_slide_count)

    # 生成每页内容
    slide_contents = {}
    for slide_info in ppt_structure.get('slides', []):
        slide_content = generate_slide_content(slide_info, topic, full_content)
        slide_contents[slide_info['slide_index']] = slide_content
        time.sleep(0.5)  # 避免API调用过于频繁

    # 生成输出路径
    # output_filename = f"{topic.replace(' ', '_')}_presentation.pptx"
    # output_path = os.path.join(os.getcwd(), output_filename)
    output_path = document_path

    # 替换PPT内容
    success = replace_ppt_content(selected_template_path, ppt_structure, slide_contents, output_path)

    if success:
        # 保存内容记录
        # record_data = {
        #     'topic': topic,
        #     'selected_template': selected_template_path,
        #     'input_analysis': input_analysis,
        #     'structure': ppt_structure,
        #     'slide_contents': slide_contents,
        #     'generated_at': time.strftime("%Y-%m-%d %H:%M:%S")
        # }

        # record_path = os.path.join(os.getcwd(), f"{topic.replace(' ', '_')}_content.json")
        # with open(record_path, 'w', encoding='utf-8') as f:
        #     json.dump(record_data, f, indent=2, ensure_ascii=False)

        return f"PPT已生成！\n主题: {topic}\n使用模板: {selected_template_path}\n保存路径为: {output_path}\n"
    else:
        return f"PPT生成失败，请检查模板文件和输入内容。"



@tool
@log_io
def ppt_generator(
        user_input: Annotated[str, "User's request/requirements"], path: Annotated[str, "path to save ppt"]
) -> HumanMessage:
    """
    使用此工具根据用户输入生成PPT。

    Args:
        user_input: 用户的PPT需求、主题、大纲或文档路径
        document_path: 文档路径（可选）

    Returns:
        生成结果信息
    """
    path = str(resolve_file_path(path))

    logger.info(f"{ppt_state}-=-PPT专员===处理需求：{user_input}")
    try:
        result = ppt_generation_tool(user_input, path)

        return HumanMessage(content=result)
    except Exception as e:
        error_msg = f"PPT生成失败: {repr(e)}"
        return HumanMessage(content=error_msg)
    
